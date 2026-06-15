"""附件解析服务（第一阶段）。

把用户在对话中上传的资料（docx/pdf/xlsx/csv/txt/md/pptx/图片）解析为可用于项目档案
补全的文本，让"附件真正进入分析"。

红线：
- 仅解析"用户自己上传的项目资料"，不读取训练语料原文、不外发；
- 文本提取后仅做长度截断与项目要素抽取，不把整篇内容写入前台/日志；
- 不支持的格式（如未配置 OCR 的图片）→ 明确返回"已接收，暂未解析内容"，绝不编造。
"""

from __future__ import annotations

import csv
import io
import logging
from typing import Any

import chardet

logger = logging.getLogger("cityrenew.attachment")

TEXT_EXTS = {"txt", "md", "csv", "tsv", "json", "log"}
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "webp"}
MAX_TEXT_CHARS = 20000   # 单文件提取文本上限，避免超大文件占用内存


def _ext(filename: str) -> str:
    return (filename.rsplit(".", 1)[-1].lower() if "." in filename else "").strip()


def _decode(data: bytes) -> str:
    try:
        guess = chardet.detect(data[:4096]) or {}
        enc = guess.get("encoding") or "utf-8"
        return data.decode(enc, errors="replace")
    except Exception:  # noqa: BLE001
        return data.decode("utf-8", errors="replace")


def _parse_docx(data: bytes) -> str:
    import docx  # python-docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            txt = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            txt = ""
        if txt.strip():
            parts.append(txt.strip())
    return "\n".join(parts)


def _parse_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _parse_csv(data: bytes) -> str:
    text = _decode(data)
    parts = []
    try:
        reader = csv.reader(io.StringIO(text))
        for row in reader:
            cells = [c.strip() for c in row if c and c.strip()]
            if cells:
                parts.append(" | ".join(cells))
    except Exception:  # noqa: BLE001
        return text[:MAX_TEXT_CHARS]
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_parts.append(t)
        if slide_parts:
            parts.append(f"[第{i}页] " + " ".join(slide_parts))
    return "\n".join(parts)


def parse(filename: str, data: bytes) -> dict[str, Any]:
    """解析单个附件，返回 {ok, filename, ext, chars, text, summary, note}。"""
    ext = _ext(filename)
    result: dict[str, Any] = {"ok": False, "filename": filename, "ext": ext,
                              "chars": 0, "text": "", "summary": "", "note": ""}
    try:
        if ext == "docx":
            text = _parse_docx(data)
        elif ext == "pdf":
            text = _parse_pdf(data)
        elif ext in ("xlsx", "xlsm"):
            text = _parse_xlsx(data)
        elif ext in ("csv", "tsv"):
            text = _parse_csv(data)
        elif ext == "pptx":
            text = _parse_pptx(data)
        elif ext in TEXT_EXTS:
            text = _decode(data)
        elif ext in IMAGE_EXTS:
            result["ok"] = True
            result["note"] = "图片已接收（暂未做图内文字识别），将作为项目资料留存。"
            result["summary"] = "图片资料"
            return result
        elif ext == "doc":
            result["note"] = "旧版 .doc 暂不支持解析，请另存为 .docx 后上传。"
            return result
        else:
            result["note"] = f"暂不支持解析的格式：.{ext}"
            return result
    except Exception as exc:  # noqa: BLE001
        logger.warning("附件解析失败 ext=%s err=%s", ext, type(exc).__name__)
        result["note"] = "文件解析失败，请确认文件未损坏或更换格式。"
        return result

    text = (text or "").strip()
    if not text:
        result["ok"] = True
        result["note"] = "未从文件中提取到可用文本。"
        return result
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS]
    result["ok"] = True
    result["chars"] = len(text)
    result["text"] = text
    result["summary"] = text[:80].replace("\n", " ")
    return result
